import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from src.modules.report_generator.schemas import ReportResponse


def _add_slide(prs: Presentation, title: str, content: str | list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    if isinstance(content, str):
        body.text = content
    else:
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(content):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def build_pptx_report(report: ReportResponse, mode: str, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide using layout with title placeholder
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = report.title or "Creative Pre-Test Report"
    mode_label = "Internal Report" if mode == "internal" else "Client-Facing Report"
    date_str = report.created_at.strftime('%Y-%m-%d') if report.created_at else 'N/A'
    from pptx.util import Inches as In
    txBox = title_slide.shapes.add_textbox(In(1), In(2), In(11), In(1))
    tf = txBox.text_frame
    tf.text = f"{mode_label}\nGenerated: {date_str}"

    _add_slide(prs, "Creative Overview", f"Test Run ID: {report.test_run_id}")
    _add_slide(prs, "Executive Summary", report.summary or "No summary available.")

    if report.scorecard:
        lines = [f"{s.criterion}: {s.score:.0f}/10" for s in report.scorecard]
        avg = sum(s.score for s in report.scorecard) / max(len(report.scorecard), 1)
        lines.append(f"")
        lines.append(f"Average score: {avg:.1f}/10")
        _add_slide(prs, "Scorecard", lines)

    if report.audience_reactions:
        lines = []
        for a in report.audience_reactions:
            lines.append(f"{a.segment_name or a.audience_profile_id}: {a.reaction}")
            lines.append(f"  Engagement: {a.engagement_probability:.0%}")
        _add_slide(prs, "Audience Reactions", lines)

    if report.risks:
        lines = [f"[{r.severity.upper()}] {r.risk}" for r in report.risks]
        _add_slide(prs, "Brand Safety Risks", lines)
    else:
        _add_slide(prs, "Brand Safety Risks", ["No significant risks identified."])

    if report.visual_notes:
        _add_slide(prs, "Visual Analysis Notes", [report.visual_notes[:500]])

    if report.recommendations:
        lines = [f"[{r.priority.upper()}] {r.recommendation}" for r in report.recommendations]
        _add_slide(prs, "Recommendations", lines)

    rec = report.final_recommendation.upper() if report.final_recommendation else "REVISE"
    _add_slide(prs, "Final Recommendation", [f"Recommendation: {rec}"])

    prs.save(output_path)
    return output_path
