"""Tests for PPTX export generation."""

import os

from httpx import ASGITransport, AsyncClient
from src.main import app

from src.modules.export_jobs.renderers.pptx_report import build_pptx_report
from src.modules.report_generator.service import generate_report
from src.modules.creative_assets.service import create_asset as create_creative_asset
from src.modules.creative_assets.schemas import CreateCreativeAssetRequest
from src.modules.test_runs.schemas import CreateTestRunRequest
from src.modules.test_runs.service import create_test_run, run_test_run
from src.shared.config.settings import get_settings


def test_pptx_renderer_creates_file():
    asset = create_creative_asset(CreateCreativeAssetRequest(
        title="PPTX Test Asset", asset_type="text", text_content="Test content for PPTX."
    ))
    run = create_test_run(CreateTestRunRequest(creative_asset_id=asset.id))
    run = run_test_run(run.id)
    report = generate_report(run.id, report_mode="internal", report_format="json")
    assert report is not None

    settings = get_settings()
    output_path = os.path.join(settings.exports_root, "test_pptx.pptx")
    os.makedirs(settings.exports_root, exist_ok=True)

    result = build_pptx_report(report, "internal", output_path)
    assert os.path.isfile(result)
    assert os.path.getsize(result) > 0
    os.remove(result)


def test_pptx_export_endpoint_creates_job():
    transport = ASGITransport(app=app)
    import anyio

    async def test():
        asset = create_creative_asset(CreateCreativeAssetRequest(
            title="PPTX Endpoint Asset", asset_type="text", text_content="PPTX endpoint test."
        ))
        run = create_test_run(CreateTestRunRequest(creative_asset_id=asset.id))
        run = run_test_run(run.id)
        report = generate_report(run.id, report_mode="internal", report_format="json")
        assert report is not None

        async with AsyncClient(transport=transport, base_url="http://test") as client:
            resp = await client.post(f"/exports/report/{report.id}/pptx")
        assert resp.status_code == 201
        data = resp.json()
        assert data["export_type"] == "pptx"
        assert data["status"] == "completed"
        assert data["file_path"] is not None

    anyio.run(test)


def test_pptx_file_exists_and_nonzero():
    transport = ASGITransport(app=app)
    import anyio

    async def test():
        asset = create_creative_asset(CreateCreativeAssetRequest(
            title="PPTX File Test", asset_type="text", text_content="Test content."
        ))
        run = create_test_run(CreateTestRunRequest(creative_asset_id=asset.id))
        run = run_test_run(run.id)
        report = generate_report(run.id, report_mode="internal", report_format="json")

        async with AsyncClient(transport=transport, base_url="http://test") as client:
            resp = await client.post(f"/exports/report/{report.id}/pptx")
        assert resp.status_code == 201
        data = resp.json()
        assert os.path.isfile(data["file_path"])
        assert os.path.getsize(data["file_path"]) > 100

    anyio.run(test)


def test_pptx_has_title_slide():
    asset = create_creative_asset(CreateCreativeAssetRequest(
        title="PPTX Title Test", asset_type="text", text_content="Check title slide."
    ))
    run = create_test_run(CreateTestRunRequest(creative_asset_id=asset.id))
    run = run_test_run(run.id)
    report = generate_report(run.id, report_mode="internal", report_format="json")
    assert report is not None

    settings = get_settings()
    output_path = os.path.join(settings.exports_root, "test_pptx_title.pptx")
    os.makedirs(settings.exports_root, exist_ok=True)
    build_pptx_report(report, "internal", output_path)

    from pptx import Presentation
    prs = Presentation(output_path)
    assert len(prs.slides) >= 1
    assert prs.slides[0].shapes.title is not None
    os.remove(output_path)


def test_pptx_expected_slide_count():
    asset = create_creative_asset(CreateCreativeAssetRequest(
        title="PPTX Count Test", asset_type="text", text_content="Count slides."
    ))
    run = create_test_run(CreateTestRunRequest(creative_asset_id=asset.id))
    run = run_test_run(run.id)
    report = generate_report(run.id, report_mode="internal", report_format="json")
    assert report is not None

    settings = get_settings()
    output_path = os.path.join(settings.exports_root, "test_pptx_count.pptx")
    os.makedirs(settings.exports_root, exist_ok=True)
    build_pptx_report(report, "internal", output_path)

    from pptx import Presentation
    prs = Presentation(output_path)
    assert 3 <= len(prs.slides) <= 12
    os.remove(output_path)
